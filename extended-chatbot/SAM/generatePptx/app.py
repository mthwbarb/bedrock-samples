
from pptx import Presentation
import json
import boto3
import os
from json import JSONDecoder

def generatePayload(bucket,uuid):
    print("getting chat history: "+ 'chats/'+uuid+'.json')
    chatObject = bucket.Object('chats/'+uuid+'.json').get()['Body'].read().decode('utf-8')
    chatObject_json = str(json.loads(chatObject))
    payload = """\n\nHuman:Given the following conversation about AWS services provided in json format, 
    summaraize the key points in the form of a presentation in json format. Use the template provided below to format the output of each slide.  Include at least 3 slides.
<conversation>= """+chatObject_json+""".  Use the following format for the output:
{
  "title": "Slide Title", 
  "bullets": [
    "Bullet 1",
    "Bullet 2", 
    "Bullet 3",
    "Bullet 4"
  ]
}
\n\nAssistant:"""
    return payload

def extract_json_objects(text, decoder=JSONDecoder()):
    """Find JSON objects in text, and yield the decoded JSON data

    Does not attempt to look for JSON arrays, text, or other JSON types outside
    of a parent JSON objesct.

    """
    pos = 0
    while True:
        match = text.find('{', pos)
        if match == -1:
            break
        try:
            result, index = decoder.raw_decode(text[match:])
            yield result
            pos = match + index
        except ValueError:
            pos = match + 1

def generateSlides(payload,client,modelId,uuid,s3Client,bucket):
    try:
        print("Invoking Endpoint")
        body = json.dumps({"prompt": payload, "max_tokens_to_sample": 2048, "temperature": 1, "top_k": 250, "top_p":0.999, "stop_sequences": ["\\n\\nHuman:","\\n\\nhuman:"]})
        accept = "application/json"
        contentType = "application/json"
        try:
            response = client.invoke_model(body=body, modelId=modelId, accept=accept, contentType=contentType)
            response_body = json.loads(response.get("body").read())
            output = response_body.get("completion")
            print(output)
        except (Exception, KeyboardInterrupt) as e:
            print(e)
            return 'Error occurred:'
        slides = []
        for result in extract_json_objects(output):
            slides.append(result)
        prs = Presentation()
        #Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Here is your Presentation"
        subtitle.text = "-Team Atlantis"
        prs.save('/tmp/output.pptx')
        #Slides
        for slide2 in slides:
            print(slide2['title'])
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
    
            title_shape.text = slide2['title']
            
            for bullet in slide2['bullets']:
                tf = body_shape.text_frame
                p = tf.add_paragraph()
                p.text = bullet
                prs.save('/tmp/output.pptx')
        print('Uploading pptx file to s3: ' + 'presentations/'+uuid+'.pptx')
        s3=s3Client
        s3.meta.client.upload_file('/tmp/output.pptx',bucket,'presentations/'+uuid+'.pptx')
        #cleanup
        print("Cleaning Up temp file")
        os.remove("/tmp/output.pptx")
        return 'presentations/'+uuid+'.pptx'
    except (Exception, KeyboardInterrupt) as e:
        print(e)
        return "Error creating pptx"
def lambda_handler(event, context):
    try:
        print("Starting...")
        print(event)
    
        bucketname = event['bucket']
        uuid = event['uuid']
        bedrockRegion = 'us-east-1'
        modelId = event['modelId']
        
    
        #print("runId: " +runId)
        print("Bucket: "+bucketname)
        bedrock = boto3.client('bedrock-runtime', bedrockRegion)
        s3=boto3.resource('s3')
        bucket = s3.Bucket(bucketname)

        print("Generating Payload for LLM")
        payload = generatePayload(bucket,uuid)
        print("Sending Payload to LLM and generating the slide deck:" + payload)
        pptxObject = generateSlides(payload,bedrock,modelId,uuid,s3,bucketname)
        return pptxObject
    except (Exception, KeyboardInterrupt) as e:
        print(e)
        return 'Error occurred:'
